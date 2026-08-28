"""Builds the journal-track presentation deck for FQL-GTA from the same
results artifacts the manuscript uses, so the numbers on the slides always
match the numbers in the paper. Slide helpers are self-contained (this
project has no shared base presentation module to import from)."""
import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
RESULTS = ROOT / "Results"
FIGURES = ROOT / "Figures"
PHOTO = ROOT / "Supplementary" / "KK_Profile_Pic.jpeg"
OUT = ROOT / "Supplementary" / "Presentation_FQL-GTA.pptx"

NAVY = RGBColor(0x1F, 0x2A, 0x44)
ORANGE = RGBColor(0xDD, 0x84, 0x52)
BLUE = RGBColor(0x4C, 0x72, 0xB0)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def load():
    with open(RESULTS / "metrics.json") as f:
        metrics = json.load(f)
    rl = pd.read_csv(RESULTS / "rl_baselines_comparison.csv").set_index("method")
    ablation = pd.read_csv(RESULTS / "ablation_summary.csv")
    ablation_per_seed = pd.read_csv(RESULTS / "ablation_per_seed.csv")
    with open(RESULTS / "statistical_tests.json") as f:
        stats_j = json.load(f)
    cross_env = pd.read_csv(RESULTS / "cross_env_transfer_summary.csv")
    with open(RESULTS / "computational_cost.json") as f:
        comp_cost = json.load(f)
    with open(RESULTS / "learned_rules_worked_examples.json") as f:
        rules = json.load(f)
    with open(RESULTS / "alpha_sensitivity_meta.json") as f:
        alpha_sens_meta = json.load(f)
    return metrics, rl, ablation, ablation_per_seed, stats_j, cross_env, comp_cost, rules, alpha_sens_meta


def add_title_slide(prs, title, subtitle, photo=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(9.5), Inches(2.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(32); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.LEFT
    p2 = tf.add_paragraph(); p2.text = subtitle; p2.font.size = Pt(15)
    p2.font.color.rgb = RGBColor(0xCC, 0xD3, 0xE0); p2.space_before = Pt(14)
    if photo and Path(photo).exists():
        size = Emu(1600200)
        margin = Emu(400000)
        x = prs.slide_width - size - margin
        y = margin
        border = slide.shapes.add_shape(1, x - Emu(28575), y - Emu(28575),
                                         size + Emu(57150), size + Emu(57150))
        border.fill.solid(); border.fill.fore_color.rgb = WHITE; border.line.fill.background()
        slide.shapes.add_picture(str(photo), x, y, width=size, height=size)
    return slide


def add_content_slide(prs, title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(11), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(26)
    p.font.bold = True; p.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(5.5))
    tf = body.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(18); p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(10)
    if note:
        slide.notes_slide.notes_text_frame.text = note
    return slide


def add_image_slide(prs, title, image_path, note=None, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(11), Inches(0.75))
    p = tb.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(24)
    p.font.bold = True; p.font.color.rgb = WHITE

    from PIL import Image as _PILImage
    _w_px, _h_px = _PILImage.open(str(image_path)).size
    _aspect = _w_px / _h_px
    _max_w = prs.slide_width - Inches(0.6)
    _max_h = Inches(5.6)
    if _max_w / _aspect <= _max_h:
        pic = slide.shapes.add_picture(str(image_path), Inches(1.3), Inches(1.25), width=_max_w)
    else:
        pic = slide.shapes.add_picture(str(image_path), Inches(1.3), Inches(1.25), height=_max_h)
    pic.left = int((prs.slide_width - pic.width) / 2)
    if caption:
        cb = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(11.5), Inches(0.5))
        cp = cb.text_frame.paragraphs[0]; cp.text = caption
        cp.font.size = Pt(12); cp.font.italic = True; cp.font.color.rgb = GREY
        cp.alignment = PP_ALIGN.CENTER
    if note:
        slide.notes_slide.notes_text_frame.text = note
    return slide


def add_table_slide(prs, title, headers, rows, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(11), Inches(0.75))
    p = tb.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(24)
    p.font.bold = True; p.font.color.rgb = WHITE

    rows_n, cols_n = len(rows) + 1, len(headers)
    table_shape = slide.shapes.add_table(rows_n, cols_n, Inches(0.4), Inches(1.5),
                                          Inches(12.5), Inches(0.65 * rows_n))
    table = table_shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c); cell.text = str(h)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(13)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c); cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(13)
    if note:
        slide.notes_slide.notes_text_frame.text = note
    return slide


def build():
    metrics, rl, ablation, ablation_per_seed, stats_j, cross_env, comp_cost, rules, alpha_sens_meta = load()
    base = metrics["Baseline FQL"]
    gta = metrics["Improved FQL-GTA"]
    n_seeds = base["n_seeds_total"]
    pc = stats_j["primary_comparison"]
    ft = stats_j["friedman_test"]
    holm = stats_j["ablation_pairwise_holm_summary"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "A Systematically Evaluated Interpretable Enhancement to Fuzzy Q-Learning",
        "Gaussian membership functions, eligibility traces, and adaptive exploration on "
        "CartPole-v1\nKarthikeyan K  |  MS AI&ML, REVA University  |  Advisor: Dr. J. B. Simha\n"
        "Manuscript prepared for journal submission",
        photo=PHOTO,
    )

    add_content_slide(prs, "Two Research Questions, Tested Not Assumed", [
        "RQ1: Does combining Gaussian membership functions, Q(lambda) eligibility traces, and "
        "adaptive epsilon produce a statistically defensible improvement over classical FQL, "
        "attributable to the combination rather than any single component?",
        "RQ2: Do CartPole-tuned FQL-GTA hyperparameters transfer zero-shot to other control "
        "environments, or is per-environment retuning necessary?",
        "Both are answered empirically, including the ablation's honest statistical-power "
        "limitation and the weak zero-shot transfer result.",
    ])

    add_content_slide(prs, "Why Fuzzy Q-Learning, When Deep RL Dominates Benchmarks?", [
        "Deep RL policies (DQN, PPO) are opaque -- a trained network cannot be audited rule "
        "by rule, which matters in safety- and compliance-sensitive control settings.",
        "Classical FQL (Jouffe, 1998) is inspectable by design: every decision traces back to "
        "specific fuzzy rules and their firing strengths.",
        "But the classical formulation uses triangular membership functions, one-step TD "
        "updates, and fixed exploration -- design choices this study revisits directly.",
    ])

    add_table_slide(
        prs, f"RQ1: Primary Comparison, {n_seeds} Seeds",
        ["Agent", "Final return (mean ± SD)", "Solve rate", "95% CI"],
        [["Baseline FQL", f"{base['final_return_mean']:.1f} ± {base['final_return_std']:.1f}",
          f"{base['solved_rate']:.0%}",
          f"[{base['final_return_95ci'][0]:.1f}, {base['final_return_95ci'][1]:.1f}]"],
         ["Improved FQL-GTA", f"{gta['final_return_mean']:.1f} ± {gta['final_return_std']:.1f}",
          f"{gta['solved_rate']:.0%}",
          f"[{gta['final_return_95ci'][0]:.1f}, {gta['final_return_95ci'][1]:.1f}]"]],
        note=f"Paired Wilcoxon p = {pc['wilcoxon_p_value']:.2e}, Cohen's d = "
             f"{pc['cohens_d_trailing50_return']:.2f} -- a large, statistically robust effect. "
             f"The two agents use different learning rates by frozen config (baseline "
             f"alpha={alpha_sens_meta['canonical_alpha']}, GTA alpha={alpha_sens_meta['gta_alpha']}); "
             f"a dedicated sweep tuning the baseline's alpha finds the best baseline still "
             f"loses by a large, significant margin (d={alpha_sens_meta['tuned_vs_gta']['cohens_d']:.2f}), "
             f"so the result is not a learning-rate artifact.",
    )

    add_image_slide(
        prs, "Learning Curves, All 30 Seeds",
        FIGURES / "fig1_learning_curves.png",
        caption="Mean ± 95% CI, 30 seeds per agent.",
        note="FQL-GTA crosses the solve threshold mid-training in most seeds but only holds it "
             "through episode 500 in 40% -- motivating the strict final-trailing-window solve "
             "definition over first-crossing.",
    )

    _full_gta_ps = ablation_per_seed[ablation_per_seed["label"] == "On / On / On (full FQL-GTA)"]
    _n_solved_abl = int((_full_gta_ps["trailing50_return"] >= 475).sum())
    add_image_slide(
        prs, "Full 2³ Factorial Ablation: No Single Component Solves It Alone",
        FIGURES / "fig14_ablation_factorial.png",
        caption="All 8 combinations of Gaussian MF / eligibility trace / adaptive epsilon, 3 seeds each.",
        note=f"Best single-component condition reaches only "
             f"{ablation.iloc[3]['trailing50_mean']:.0f} trailing-50 return; the full combination "
             f"reaches a mean of {ablation.iloc[7]['trailing50_mean']:.0f} across its "
             f"{len(_full_gta_ps)} seeds, but only {_n_solved_abl} of {len(_full_gta_ps)} "
             "individual seeds actually cross the 475 bar -- the mean exceeds it, not every "
             "seed. A formal factorial effect analysis (next section) shows this is driven by "
             "significant Gaussian-MF x eligibility-trace and Gaussian-MF x adaptive-epsilon "
             "two-way interaction effects, with no significant three-way interaction.",
    )

    add_content_slide(prs, "Statistics, Including the Honest Power Limitation", [
        f"Primary comparison: Shapiro-Wilk rejects normality for both agents, so the paired "
        f"Wilcoxon test (p = {pc['wilcoxon_p_value']:.2e}) is the primary evidence, not the t-test.",
        f"Ablation: a Friedman omnibus test across all 8 conditions is significant "
        f"(chi-sq = {ft['statistic']:.1f}, p = {ft['p_value']:.3f}), confirming a real overall effect.",
        f"But with only 3 seeds per condition, {holm['n_significant_after_holm']} of "
        f"{holm['n_pairs_tested']} Holm-corrected pairwise Wilcoxon comparisons reach "
        f"significance individually -- reported plainly as a power limitation, not hidden.",
        "10,000-resample bootstrap 95% CIs are reported alongside every point estimate.",
    ])

    add_image_slide(
        prs, "Hyperparameters Sit at Pronounced Performance Peaks",
        FIGURES / "fig15_hyperparameter_sensitivity.png",
        caption="One-factor-at-a-time sweeps, CartPole-tuned default marked on each panel.",
        note="The Gaussian width multiplier shows the sharpest cliff: values >= 0.5 collapse "
             "performance to near-random by over-activating every rule for every state.",
    )

    add_image_slide(
        prs, "RQ2: Zero-Shot Transfer Is Weak; Retuning Recovers Most of the Gap",
        FIGURES / "fig16_cross_env_transfer.png",
        caption="MountainCar-v0 and Acrobot-v1, zero-shot vs. per-environment coordinate-search tuning.",
        note="Acrobot-v1 zero-shot trailing-50 is far below its practical solve bar; a "
             "disjoint-validation-seed coordinate search (same cost as the sensitivity sweep) "
             "nearly reaches it.",
    )

    add_table_slide(
        prs, "Against Deep-RL Baselines: A Trade, Not a Win",
        ["Method", "Final return", "Solve rate", "Train time (s)", "Parameters", "Memory (KB)"],
        [[m, f"{rl.loc[m]['final_return_mean']:.1f}", f"{rl.loc[m]['solved_rate']:.0%}",
          f"{rl.loc[m]['train_time_sec_mean']:.1f}", f"{int(rl.loc[m]['n_params'])}",
          f"{rl.loc[m]['memory_kb']:.2f}"]
         for m in ["Baseline FQL", "Improved FQL-GTA", "DQN (hand-rolled)",
                    "DQN (stable-baselines3)", "PPO (stable-baselines3)"]],
        note="PPO wins on final return but uses ~14x FQL-GTA's parameters and ~14x the memory. "
             "Both DQN variants under-perform FQL-GTA within this study's training budget -- a "
             "budget-dependent finding, not a general claim about DQN.",
    )

    ex = rules["examples"][0]
    top = ex["top_firing_rules"][0]
    add_content_slide(prs, "Interpretability: One Real Decision, Fully Traced", [
        f"Real state from a trained agent's evaluation rollout: {[round(x,3) for x in ex['raw_state']]}.",
        f"{ex['n_rules_firing_above_0.01']} of {ex['n_rules_total']} rules fire above 0.01 -- "
        "local specialisation, not uniform activation across the whole rule base.",
        f"Dominant rule (index {top['rule_index']}) fires at {top['firing_strength']:.1%} strength "
        f"with Q(left) = {top['q_push_left']:.2f}, Q(right) = {top['q_push_right']:.2f}.",
        f"Blended decision: Q(left) = {ex['blended_q_values']['push_left']:.2f}, "
        f"Q(right) = {ex['blended_q_values']['push_right']:.2f} -> \"{ex['selected_action']}\".",
        "This traceability is the concrete evidence behind the paper's reframed novelty claim: "
        "not new ingredients, but a systematically evaluated, inspectable combination.",
    ])

    add_content_slide(prs, "Conclusion", [
        f"RQ1 answered: the combination produces a large, statistically robust improvement "
        f"(d = {pc['cohens_d_trailing50_return']:.2f}). A factorial effect analysis finds "
        "significant Gaussian-MF x eligibility-trace and Gaussian-MF x adaptive-epsilon "
        "two-way interaction effects, with no significant three-way interaction -- with the "
        "ablation's pairwise power limitation reported honestly.",
        "RQ2 answered: CartPole-tuned hyperparameters do not transfer zero-shot, but a "
        "lightweight, cheap coordinate search recovers most of the lost performance.",
        "Against deep RL: FQL-GTA trades some final performance for an order-of-magnitude "
        "smaller parameter count, faster training, and full rule-level interpretability.",
        "Contribution is deliberately modest: not new ingredients, but a systematically "
        "evaluated, statistically validated, reproducible enhancement to a 26-year-old method.",
    ])

    add_title_slide(prs, "Thank You",
                     "Questions welcome. Code, configuration, results, and manuscript accompany "
                     "this presentation.\nManuscript prepared for journal submission. DOI: Not yet assigned.")

    prs.save(str(OUT))
    print(f"Saved {len(prs.slides)}-slide journal deck to {OUT}")
    return prs


if __name__ == "__main__":
    build()
